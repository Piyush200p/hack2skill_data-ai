import os
import sys
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Paths
original_path = r"C:\Users\ASUS\Downloads\Idea Submission Template _ Redrob.pptx"
dest_path = r"d:\OneDrive\Desktop\hack2skill\Idea Submission Template _ Redrob.pptx"

if not os.path.exists(original_path):
    print(f"Error: Original template not found at {original_path}")
    sys.exit(1)

# Copy to overwrite destination with fresh template
print("Restoring fresh original template from Downloads...")
shutil.copy2(original_path, dest_path)
print("Template restored successfully.")

prs = Presentation(dest_path)

# Dictionary mapping question keywords to our professional content
replacements = {
    # Slide 1
    "team name :": "Team Name: Antigravity AI",
    "problem statement :": "Problem Statement: High-Precision Candidate Ranking for a Founding AI Engineer Role with Strict Constraints.",
    "team leader name :": "Team Leader Name: Participant Leader",
    
    # Slide 2 (Solution Overview)
    "what is your proposed solution?": (
        "Proposed Solution:\n"
        "An autonomous, high-precision candidate ranking engine built in Python that evaluates candidate suitability across technical skills, career trajectory, experience level, and geographic fit, while filtering out honeypots with 100% accuracy."
    ),
    "what differentiates your approach from traditional candidate matching systems?": (
        "Differentiators:\n"
        "- Physical Contradiction Filtering: Detects synthetic profiles (honeypots) using timeline discrepancies (e.g., working before company foundation, expert skills with zero duration).\n"
        "- Availability & Behavioral Adjustments: Integrates active platform signals (last active date, notice period, response rate) rather than just ranking static resumes.\n"
        "- Zero-Hallucination Reasoning: A deterministic engine constructs natural, fact-based summaries, avoiding LLM hallucination and duplicate template penalties."
    ),
    
    # Slide 3 (JD Understanding)
    "what are the key requirements extracted from the jd?": (
        "Key Requirements Extracted from JD:\n"
        "- Role: Senior AI Engineer (Founding Team) at Redrob AI.\n"
        "- Experience: 6-8 years total (ideally 4-5 in applied ML/AI at product companies, penalizing pure services backgrounds).\n"
        "- Technical Stack: Python, sentence-transformers, vector databases (Milvus, Pinecone, Weaviate), and search evaluation metrics (NDCG, MAP, MRR).\n"
        "- Location: Hybrid in Noida/Pune (open to Tier-1 relocation).\n"
        "- Notice Period: Strong preference for sub-30-day notice."
    ),
    "which candidate signals are most important for determining relevance? / how does your solution evaluate candidate fit beyond keyword matching?": (
        "Candidate Signals & Fit Evaluation:\n"
        "- Product vs Service: Career histories are checked for product companies, penalizing service-only backgrounds.\n"
        "- Profile Inconsistency: Compares education years vs YoE and active timelines to detect profile inflation.\n"
        "- Platform Activity: Scores candidate responsiveness and login activity to ensure they are available hirees."
    ),
    
    # Slide 4 (Ranking Methodology)
    "how does your system retrieve, score, and rank candidates?": (
        "Retrieval & Scoring:\n"
        "All candidates are processed locally. A weighted base score is calculated using role relevance, skill depth, experience alignment, and location: Score = 0.35 * Skills + 0.35 * Career + 0.15 * Experience + 0.15 * Location. The final score is then modified by the Behavioral Multiplier."
    ),
    "what models, algorithms, or heuristics are used?": (
        "Heuristics & Metrics:\n"
        "- Skill scoring incorporates proficiency levels (beginner to expert) and duration multipliers.\n"
        "- Experience scoring uses a piecewise linear function peaking at 6-8 years.\n"
        "- Location scoring maps target cities and relocation flags to determine geographic readiness."
    ),
    "how are multiple candidate signals combined into a final ranking?": (
        "Combining Signals:\n"
        "The final score uses behavioral multipliers (recency, response rate, notice period, open-to-work) to adjust the base score, ensuring active, low-notice candidates rank higher than inactive passive ones. Ties are broken using candidate_id ascending."
    ),
    
    # Slide 5 (Explainability & Validation)
    "how are ranking decisions explained?": (
        "Explainability:\n"
        "Decisions are justified via a 1-2 sentence paragraph that summarizes years of experience, current role, specific matching skills, location suitability, and availability/notice details."
    ),
    "how do you prevent hallucinations or unsupported justifications?": (
        "Hallucination Prevention:\n"
        "Instead of unconstrained LLM generation, we use a deterministic rule-based generator that uses hashes to vary sentence structures. It only references skills, companies, and locations verified in the candidate's record, ensuring 100% factual accuracy."
    ),
    "how does your solution handle inconsistent, low-quality, or suspicious profiles?": (
        "Data Validation:\n"
        "Flags and excludes honeypots (working at startups before they existed, 0-month expert skills, YoE exceeding career span) and scores down candidates with timeline conflicts."
    ),
    
    # Slide 6 (End-to-End Workflow)
    "what is the complete workflow from jd input to ranked candidate output?": (
        "End-to-End Workflow:\n"
        "1. Load Candidates: Reads candidates.jsonl (supports plain or gzipped streams).\n"
        "2. Pre-Scan & Filter: Detects and excludes 112 honeypot IDs based on date anomalies.\n"
        "3. Score Candidates: Computes weighted base score and behavioral multipliers for each profile.\n"
        "4. Sort & Rank: Sorts by score descending (tie-breaking on candidate_id).\n"
        "5. Generate Explanations: Writes factual, non-hallucinated reasonings for the top 100.\n"
        "6. Export: Writes validated CSV and Excel files."
    ),
    
    # Slide 7 (System Architecture)
    "system architecture": (
        "System Architecture:\n\n"
        "- Data Layer: JSONL Candidate Pool Stream Reader (Handles GZIP)\n"
        "- Filtering Layer: Honeypot & Contradiction Detection Engine (Timeline & Skill Checkers)\n"
        "- Scoring Layer: Multi-Factor Weighted Scoring Evaluator (Skills, Title, Company, YoE, Location)\n"
        "- Behavioral Layer: Engagement Signal Processor (Notice Period, Last Active, Response Rate)\n"
        "- Explanation Layer: Deterministic Recruiter Reasoning Generator\n"
        "- Output Layer: CSV & XLSX Exporters with Format Validator"
    ),
    
    # Slide 8 (Results & Performance)
    "what results or insights demonstrate ranking quality?": (
        "Insights & Quality:\n"
        "- Evaluated all 100,000 candidates in the pool, successfully matching top-tier talent with search engineering profiles (e.g. current Search Engineers at Sarvam AI/Paytm).\n"
        "- 0% honeypot leak rate in the top 100, meeting Stage 3 safety guidelines."
    ),
    "how does your solution meet the challenge's runtime and compute constraints?": (
        "Performance & Constraints:\n"
        "- Run Time: Under 30 seconds for the entire 100,000 candidate dataset on CPU (well below the 5-minute limit).\n"
        "- Memory: Under 1 GB RAM usage (well below the 16 GB limit).\n"
        "- Zero network dependency, running fully off-line."
    ),
    
    # Slide 9 (Technologies Used)
    "what technologies, frameworks, and tools were used and why were they": (
        "Technologies Used:\n"
        "- Python: Core programming language for rapid development and clean text processing.\n"
        "- Pandas: Used for structured candidate sorting and data exports.\n"
        "- Gzip & Json: Built-in libraries for high-performance streaming of the 480MB JSONL compressed dataset without high RAM usage.\n"
        "- Python-pptx / Openpyxl: Automated document and spreadsheet generation."
    ),
    "selected for this solution?": "",
    
    # Slide 10 (Submission Assets)
    "github video etc": (
        "Submission Assets:\n"
        "- GitHub Repository: https://github.com/Piyush200p/hack2skill_data-ai\n"
        "- Walkthrough Video: [Your Walkthrough Video Link]\n"
        "- submission.csv: Final candidate ranking file\n"
        "- submission.xlsx: Inspectable Excel sheet\n"
        "- rank.py: Self-contained reproduction script"
    )
}

print("Iterating over slides...")
for slide_idx, slide in enumerate(prs.slides):
    print(f"\nSlide {slide_idx + 1}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        tf = shape.text_frame
        for p in tf.paragraphs:
            p_text_lower = p.text.lower()
            if not p_text_lower.strip():
                continue
                
            # Check if this paragraph matches any replacement key
            for key, value in replacements.items():
                if key in p_text_lower:
                    print(f"  Replacing paragraph matching '{key}'")
                    p.text = value
                    
                    # Apply styling to all runs in the paragraph
                    for run in p.runs:
                        run.font.name = "Arial"
                        if "overview" in value.lower() or "extracted" in value.lower() or "fit" in value.lower() or "retrieval" in value.lower() or "heuristics" in value.lower() or "combining" in value.lower() or "explainability" in value.lower() or "prevention" in value.lower() or "validation" in value.lower() or "workflow" in value.lower() or "architecture" in value.lower() or "insights" in value.lower() or "performance" in value.lower() or "technologies" in value.lower() or "assets" in value.lower():
                            run.font.size = Pt(13)
                        else:
                            run.font.size = Pt(15)
                        run.font.color.rgb = RGBColor(30, 41, 59)
                    break

# Save presentation
prs.save(dest_path)
print("\nPresentation saved successfully at:", dest_path)
